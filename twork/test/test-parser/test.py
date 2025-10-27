#!/usr/bin/env python3
"""
StructureFactory 测试脚本

测试不同格式文档的解析、需求分析和WBS生成功能。
"""

import os
import sys
import json
from pathlib import Path
from typing import Dict, Any

# 添加twork模块路径
sys.path.insert(0, str(Path(__file__).parent.parent.parent))

from twork.parser import StructureUnderstandFactory
from twork.llm import OpenAIAdapter, LLMConfig


def create_test_documents():
    """创建测试文档（PDF、DOCX、PPTX）"""
    print("=" * 60)
    print("正在创建测试文档...")
    print("=" * 60)
    
    data_dir = Path("./data")
    
    # 创建PDF文档
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        
        pdf_path = data_dir / "sample.pdf"
        doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # 添加内容
        story.append(Paragraph("智能家居系统需求文档", styles['Title']))
        story.append(Spacer(1, 12))
        story.append(Paragraph("项目概述", styles['Heading1']))
        story.append(Paragraph("开发一个智能家居控制系统，支持设备联动、远程控制等功能。", styles['Normal']))
        story.append(Spacer(1, 12))
        story.append(Paragraph("功能需求", styles['Heading1']))
        story.append(Paragraph("1. 设备管理 - 添加、删除、配置智能设备", styles['Normal']))
        story.append(Paragraph("2. 场景联动 - 设置自动化场景", styles['Normal']))
        story.append(Paragraph("3. 远程控制 - 手机APP远程控制设备", styles['Normal']))
        story.append(Spacer(1, 12))
        story.append(Paragraph("非功能需求", styles['Heading1']))
        story.append(Paragraph("性能：设备响应时间 < 500ms", styles['Normal']))
        story.append(Paragraph("安全：数据传输加密", styles['Normal']))
        
        doc.build(story)
        print(f"✓ 创建PDF文档: {pdf_path}")
    except ImportError:
        print("⚠ reportlab未安装，跳过PDF生成")
    except Exception as e:
        print(f"✗ PDF生成失败: {e}")
    
    # 创建DOCX文档
    try:
        from docx import Document
        
        docx_path = data_dir / "sample.docx"
        doc = Document()
        
        doc.add_heading('在线教育平台需求', 0)
        doc.add_heading('项目背景', level=1)
        doc.add_paragraph('开发一个在线教育平台，提供课程管理、直播教学、作业提交等功能。')
        
        doc.add_heading('核心功能', level=1)
        doc.add_paragraph('课程管理 - 课程创建、发布、编辑', style='List Bullet')
        doc.add_paragraph('在线学习 - 视频播放、进度跟踪', style='List Bullet')
        doc.add_paragraph('作业系统 - 作业发布、提交、批改', style='List Bullet')
        doc.add_paragraph('互动讨论 - 课程问答、论坛交流', style='List Bullet')
        
        doc.add_heading('技术要求', level=1)
        doc.add_paragraph('支持高清视频播放')
        doc.add_paragraph('支持1000+并发在线学习')
        doc.add_paragraph('移动端适配')
        
        doc.save(str(docx_path))
        print(f"✓ 创建DOCX文档: {docx_path}")
    except ImportError:
        print("⚠ python-docx未安装，跳过DOCX生成")
    except Exception as e:
        print(f"✗ DOCX生成失败: {e}")
    
    # 创建PPTX文档
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        pptx_path = data_dir / "sample.pptx"
        prs = Presentation()
        
        # 幻灯片1：标题
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title1 = slide1.shapes.title
        subtitle1 = slide1.placeholders[1]
        title1.text = "智慧物流系统"
        subtitle1.text = "项目需求说明"
        
        # 幻灯片2：项目概述
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title2 = slide2.shapes.title
        content2 = slide2.placeholders[1]
        title2.text = "项目概述"
        tf = content2.text_frame
        tf.text = "开发智慧物流管理系统"
        p = tf.add_paragraph()
        p.text = "实现订单管理、仓储管理、配送管理"
        p.level = 1
        
        # 幻灯片3：核心功能
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        title3 = slide3.shapes.title
        content3 = slide3.placeholders[1]
        title3.text = "核心功能"
        tf3 = content3.text_frame
        tf3.text = "订单管理"
        
        for item in ["仓储管理", "配送路线优化", "实时追踪", "数据分析"]:
            p = tf3.add_paragraph()
            p.text = item
            p.level = 0
        
        prs.save(str(pptx_path))
        print(f"✓ 创建PPTX文档: {pptx_path}")
    except ImportError:
        print("⚠ python-pptx未安装，跳过PPTX生成")
    except Exception as e:
        print(f"✗ PPTX生成失败: {e}")
    
    print()


def test_document_parsing():
    """测试文档解析功能"""
    print("=" * 60)
    print("测试1: 文档解析功能")
    print("=" * 60)
    
    from twork.parser.tools import DocParseTool
    
    data_dir = Path("./data")
    doc_parser = DocParseTool()
    
    # 测试所有文档格式
    test_files = [
        "sample.md",
        "sample.txt",
        "sample.pdf",
        "sample.docx",
        "sample.pptx"
    ]
    
    for filename in test_files:
        file_path = data_dir / filename
        if not file_path.exists():
            print(f"⚠ 文件不存在，跳过: {filename}")
            continue
        
        try:
            print(f"\n测试文件: {filename}")
            result = doc_parser.execute({"file_path": str(file_path)})
            sections = result.get("sections", [])
            print(f"  ✓ 解析成功，共 {len(sections)} 个章节")
            
            # 显示前3个章节
            for i, section in enumerate(sections[:3]):
                title = section.get("title", "无标题")
                content_preview = section.get("content", "")[:50]
                print(f"    - [{i+1}] {title}: {content_preview}...")
            
            if len(sections) > 3:
                print(f"    ... 还有 {len(sections) - 3} 个章节")
                
        except Exception as e:
            print(f"  ✗ 解析失败: {e}")
    
    print()


def test_structure_factory_without_llm():
    """测试StructureFactory基本功能（不使用LLM）"""
    print("=" * 60)
    print("测试2: StructureFactory基本功能（不使用LLM）")
    print("=" * 60)
    
    data_dir = Path("./data")
    test_file = data_dir / "sample.md"
    
    if not test_file.exists():
        print(f"✗ 测试文件不存在: {test_file}")
        return
    
    try:
        # 创建工厂实例
        factory = StructureUnderstandFactory(
            project_id="test-project-001",
            original_file_path=str(test_file),
            cache_dir="./cache"
        )
        
        print(f"✓ 工厂创建成功")
        print(f"  项目ID: {factory.project_id}")
        print(f"  文档路径: {factory.original_file_path}")
        print(f"  缓存目录: {factory.cache_dir}")
        print(f"  可用工具: {list(factory.tools.keys())}")
        
        # 测试文档解析（不需要LLM）
        print("\n执行文档解析...")
        parsed = factory._get_or_parse_document()
        sections = parsed.get("sections", [])
        print(f"✓ 文档解析成功，共 {len(sections)} 个章节")
        
        # 检查缓存
        cache_file = factory.cache_dir / factory.project_id / "parsed_text.json"
        if cache_file.exists():
            print(f"✓ 缓存已生成: {cache_file}")
        
    except Exception as e:
        print(f"✗ 测试失败: {e}")
        import traceback
        traceback.print_exc()
    
    print()


def test_structure_factory_with_llm():
    """测试StructureFactory完整流程（使用LLM）"""
    print("=" * 60)
    print("测试3: StructureFactory完整流程（使用LLM）")
    print("=" * 60)
    
    # 检查API配置
    api_key = os.getenv("OPENAI_API_KEY")
    api_base = os.getenv("OPENAI_API_BASE", "https://api.openai.com/v1")
    
    if not api_key:
        print("⚠ 未配置OPENAI_API_KEY环境变量，跳过LLM测试")
        print("  提示: 设置环境变量后可测试完整功能")
        print("  export OPENAI_API_KEY=your-api-key")
        return
    
    data_dir = Path("./data")
    test_file = data_dir / "sample.md"
    
    if not test_file.exists():
        print(f"✗ 测试文件不存在: {test_file}")
        return
    
    try:
        # 配置LLM
        print("配置LLM适配器...")
        llm_config = LLMConfig(
            api_base_url=api_base,
            api_key=api_key,
            model_name="qwen3-max",
            temperature=0.7
        )
        llm = OpenAIAdapter(llm_config)
        print("✓ LLM配置完成")
        
        # 创建工厂
        factory = StructureUnderstandFactory(
            project_id="test-project-002",
            original_file_path=str(test_file),
            cache_dir="./cache"
        )
        
        # 配置工具
        print("\n配置分析工具...")
        from twork.parser.templates.extraction_schema import (
            get_json_schema_for_requirement_analysis,
            get_json_schema_for_wbs
        )
        
        factory.tools["analyzer"].setup({
            "llm_adapter": llm,
            "response_format": get_json_schema_for_requirement_analysis()
        })
        factory.tools["wbs"].setup({
            "llm_adapter": llm,
            "response_format": get_json_schema_for_wbs()
        })
        print("✓ 工具配置完成")
        
        # 执行完整流程
        print("\n执行完整流程...")
        result = factory.run()
        
        # 显示结果
        print("\n" + "=" * 60)
        print("执行结果:")
        print("=" * 60)
        
        # 需求和领域分析结果
        req_domain = result.get("requirements_and_domain", {})
        domain = req_domain.get("domain", "未知")
        func_reqs = req_domain.get("functional_requirements", [])
        non_func_reqs = req_domain.get("non_functional_requirements", [])
        
        print(f"\n领域识别: {domain}")
        print(f"\n功能需求 ({len(func_reqs)} 项):")
        for req in func_reqs[:5]:
            print(f"  - [{req.get('id')}] {req.get('desc')} (优先级: {req.get('priority', '中')})")
        if len(func_reqs) > 5:
            print(f"  ... 还有 {len(func_reqs) - 5} 项")
        
        print(f"\n非功能需求 ({len(non_func_reqs)} 项):")
        for req in non_func_reqs[:5]:
            print(f"  - [{req.get('id')}] ({req.get('type')}) {req.get('desc')}")
        if len(non_func_reqs) > 5:
            print(f"  ... 还有 {len(non_func_reqs) - 5} 项")
        
        # WBS结果
        wbs = result.get("wbs", {})
        phase = wbs.get("phase", "未知")
        tasks = wbs.get("tasks", [])
        
        print(f"\nWBS阶段: {phase}")
        print(f"任务数量: {count_all_tasks(tasks)} 个")
        print(f"\n任务结构预览:")
        show_task_tree(tasks, max_level=2)
        
        # 保存结果
        output_file = Path("./output_result.json")
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print(f"\n✓ 完整结果已保存到: {output_file}")
        
    except Exception as e:
        print(f"\n✗ 测试失败: {e}")
        import traceback
        traceback.print_exc()
    
    print()


def count_all_tasks(tasks: list) -> int:
    """递归统计所有任务数量"""
    count = len(tasks)
    for task in tasks:
        children = task.get("children", [])
        count += count_all_tasks(children)
    return count


def show_task_tree(tasks: list, level: int = 0, max_level: int = 3):
    """显示任务树结构"""
    if level >= max_level:
        return
    
    for task in tasks:
        indent = "  " * level
        task_id = task.get("task_id", "")
        task_name = task.get("task_name", "")
        estimated_hours = task.get("estimated_hours", 0)
        
        print(f"{indent}- [{task_id}] {task_name} ({estimated_hours}h)")
        
        children = task.get("children", [])
        if children:
            show_task_tree(children, level + 1, max_level)


def main():
    """主测试函数"""
    print("\n" + "=" * 60)
    print("StructureFactory 测试套件")
    print("=" * 60)
    print()
    
    # 1. 创建测试文档
    create_test_documents()
    
    # 2. 测试文档解析
    test_document_parsing()
    
    # 3. 测试基本功能
    test_structure_factory_without_llm()
    
    # 4. 测试完整流程（需要LLM）
    test_structure_factory_with_llm()
    
    print("=" * 60)
    print("测试完成！")
    print("=" * 60)


if __name__ == "__main__":
    main()
