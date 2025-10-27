"""
DocParseTool - 文档解析工具

负责解析文档为带层级的结构化文本，支持PDF、DOCX、TXT、Markdown、PPTX等格式。
"""

from typing import Dict, Any, List
from pathlib import Path
import json
import PyPDF2
import docx
from pptx import Presentation
from .base_tool import BaseTool
from ...utils.logger import get_logger

logger = get_logger(__name__)


class DocParseTool(BaseTool):
    """文档解析工具"""
    
    SUPPORTED_FORMATS = [".pdf", ".docx", ".txt", ".md", ".pptx"]
    
    def execute(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        解析文档
        
        Args:
            input_data: {
                "file_path": str,  # 文档路径
            }
        
        Returns:
            {
                "sections": [
                    {
                        "title": str,
                        "content": str,
                        "level": int
                    },
                    ...
                ]
            }
        """
        file_path = Path(input_data["file_path"])
        
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        file_ext = file_path.suffix.lower()
        
        if file_ext not in self.SUPPORTED_FORMATS:
            raise ValueError(
                f"不支持的文件格式: {file_ext}。支持的格式: {', '.join(self.SUPPORTED_FORMATS)}"
            )
        
        logger.info(f"开始解析文档: {file_path.name}")
        
        # 根据文件类型选择解析方法
        if file_ext == ".pdf":
            sections = self._parse_pdf(file_path)
        elif file_ext == ".docx":
            sections = self._parse_docx(file_path)
        elif file_ext == ".pptx":
            sections = self._parse_pptx(file_path)
        else:  # .md, .txt
            sections = self._parse_text(file_path)
        
        logger.info(f"文档解析完成: {len(sections)} 个章节")
        
        return {"sections": sections}
    
    def _parse_pdf(self, file_path: Path) -> List[Dict[str, Any]]:
        """解析PDF文件"""
        sections = []
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    content = page.extract_text().strip()
                    if content:
                        sections.append({
                            "title": f"Page {i + 1}",
                            "content": content,
                            "level": 1
                        })
        except Exception as e:
            logger.error(f"PDF解析失败: {str(e)}")
            raise
        
        return sections
    
    def _parse_docx(self, file_path: Path) -> List[Dict[str, Any]]:
        """解析DOCX文件"""
        sections = []
        try:
            doc = docx.Document(file_path)
            current_section = None
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                # 判断是否为标题
                if para.style.name.startswith("Heading"):
                    # 提取标题级别
                    level = int(para.style.name.replace("Heading ", "")) if len(para.style.name) > 7 else 1
                    
                    # 保存上一个章节
                    if current_section:
                        sections.append(current_section)
                    
                    # 创建新章节
                    current_section = {
                        "title": text,
                        "content": "",
                        "level": level
                    }
                else:
                    # 添加到当前章节内容
                    if current_section:
                        current_section["content"] += text + "\n"
                    else:
                        # 如果没有章节，创建默认章节
                        current_section = {
                            "title": "Introduction",
                            "content": text + "\n",
                            "level": 1
                        }
            
            # 添加最后一个章节
            if current_section:
                sections.append(current_section)
                
        except Exception as e:
            logger.error(f"DOCX解析失败: {str(e)}")
            raise
        
        return sections
    
    def _parse_pptx(self, file_path: Path) -> List[Dict[str, Any]]:
        """解析PPTX文件"""
        sections = []
        try:
            prs = Presentation(file_path)
            
            for i, slide in enumerate(prs.slides):
                title = ""
                content = ""
                
                # 提取标题和内容
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.is_placeholder and shape.placeholder_format.type == 1:  # 标题
                            title = shape.text.strip()
                        else:
                            content += shape.text.strip() + "\n"
                
                if title or content:
                    sections.append({
                        "title": title or f"Slide {i + 1}",
                        "content": content.strip(),
                        "level": 1
                    })
                    
        except Exception as e:
            logger.error(f"PPTX解析失败: {str(e)}")
            raise
        
        return sections
    
    def _parse_text(self, file_path: Path) -> List[Dict[str, Any]]:
        """解析文本文件（TXT、Markdown）"""
        sections = []
        try:
            # 尝试UTF-8编码
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
            except UnicodeDecodeError:
                # 尝试GBK编码
                with open(file_path, "r", encoding="gbk") as f:
                    content = f.read()
            
            # 如果是Markdown，尝试按标题拆分
            if file_path.suffix.lower() == ".md":
                sections = self._split_markdown(content)
            else:
                # 普通文本，作为单个章节
                sections = [{
                    "title": file_path.stem,
                    "content": content.strip(),
                    "level": 1
                }]
                
        except Exception as e:
            logger.error(f"文本文件解析失败: {str(e)}")
            raise
        
        return sections
    
    def _split_markdown(self, content: str) -> List[Dict[str, Any]]:
        """拆分Markdown内容为章节"""
        import re
        
        sections = []
        lines = content.split("\n")
        current_section = None
        
        for line in lines:
            # 匹配Markdown标题
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
            
            if heading_match:
                # 保存上一个章节
                if current_section:
                    sections.append(current_section)
                
                # 创建新章节
                level = len(heading_match.group(1))
                title = heading_match.group(2).strip()
                current_section = {
                    "title": title,
                    "content": "",
                    "level": level
                }
            else:
                # 添加到当前章节内容
                if current_section:
                    current_section["content"] += line + "\n"
                else:
                    # 如果没有章节，创建默认章节
                    current_section = {
                        "title": "Introduction",
                        "content": line + "\n",
                        "level": 1
                    }
        
        # 添加最后一个章节
        if current_section:
            sections.append(current_section)
        
        return sections if sections else [{
            "title": "Document",
            "content": content,
            "level": 1
        }]
