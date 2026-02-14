"""文档解析器 - 解析多种格式的项目文档"""
from typing import Optional, List
import os
import logging

from ..schemas.document import Document, DocumentType, ParsedDocument, ExtractedEntity, ExtractedRelation
from .entity_extractor import EntityExtractor

logger = logging.getLogger(__name__)


class DocumentParser:
    """文档解析器
    
    支持解析多种格式的项目文档：
    - PDF
    - DOCX
    - Markdown
    - TXT
    - PPTX
    """
    
    def __init__(self, entity_extractor: EntityExtractor):
        self.entity_extractor = entity_extractor
    
    def parse(self, document: Document) -> ParsedDocument:
        """解析文档
        
        Args:
            document: 文档对象
        
        Returns:
            解析后的文档
        """
        # 1. 读取文档内容
        content = self._read_document(document)
        
        # 2. 提取实体
        parsed_doc = self.entity_extractor.extract(document.id, document.project_id, content)
        
        return parsed_doc
    
    def _read_document(self, document: Document) -> str:
        """读取文档内容"""
        file_path = document.file_path
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found: {file_path}")
        
        if document.file_type == DocumentType.TXT:
            return self._read_txt(file_path)
        elif document.file_type == DocumentType.MARKDOWN:
            return self._read_txt(file_path)
        elif document.file_type == DocumentType.PDF:
            return self._read_pdf(file_path)
        elif document.file_type == DocumentType.DOCX:
            return self._read_docx(file_path)
        elif document.file_type == DocumentType.PPTX:
            return self._read_pptx(file_path)
        else:
            return self._read_txt(file_path)
    
    def _read_txt(self, file_path: str) -> str:
        """读取文本文件"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    
    def _read_pdf(self, file_path: str) -> str:
        """读取PDF文件"""
        try:
            import pypdf
            
            text = []
            with open(file_path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text.append(page.extract_text())
            
            return "\n".join(text)
        except ImportError:
            logger.warning("pypdf not installed, cannot parse PDF")
            return ""
    
    def _read_docx(self, file_path: str) -> str:
        """读取DOCX文件"""
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs]
            
            return "\n".join(paragraphs)
        except ImportError:
            logger.warning("python-docx not installed, cannot parse DOCX")
            return ""
    
    def _read_pptx(self, file_path: str) -> str:
        """读取PPTX文件"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            texts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            
            return "\n".join(texts)
        except ImportError:
            logger.warning("python-pptx not installed, cannot parse PPTX")
            return ""
    
    @staticmethod
    def detect_file_type(filename: str) -> DocumentType:
        """检测文件类型"""
        ext = os.path.splitext(filename)[1].lower()
        
        type_map = {
            ".pdf": DocumentType.PDF,
            ".docx": DocumentType.DOCX,
            ".doc": DocumentType.DOCX,
            ".md": DocumentType.MARKDOWN,
            ".markdown": DocumentType.MARKDOWN,
            ".txt": DocumentType.TXT,
            ".pptx": DocumentType.PPTX,
            ".ppt": DocumentType.PPTX,
        }
        
        return type_map.get(ext, DocumentType.UNKNOWN)